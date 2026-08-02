"""Office 文档：docx / xlsx / pptx（现代格式；legacy .doc/.xls/.ppt 不支持）。"""
from __future__ import annotations

from pathlib import Path

from ..models import ExtractError
from .base import ExtractContext, Extractor

MAX_ROWS_PER_SHEET = 2000


class DocxExtractor(Extractor):
    name = "docx"
    exts = (".docx",)

    def extract(self, path: Path, ctx: ExtractContext) -> str:
        import docx

        try:
            doc = docx.Document(str(path))
        except Exception as e:
            raise ExtractError(f"docx 解析失败: {e}") from e
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for table in doc.tables:
            for row in table.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append("\t".join(cells))
        return "\n".join(parts)


class XlsxExtractor(Extractor):
    name = "xlsx"
    exts = (".xlsx", ".xlsm")

    def extract(self, path: Path, ctx: ExtractContext) -> str:
        import openpyxl

        try:
            wb = openpyxl.load_workbook(str(path), read_only=True, data_only=True)
        except Exception as e:
            raise ExtractError(f"xlsx 解析失败: {e}") from e
        parts = []
        try:
            for ws in wb.worksheets:
                parts.append(f"# 工作表: {ws.title}")
                for i, row in enumerate(ws.iter_rows(values_only=True)):
                    if i >= MAX_ROWS_PER_SHEET:
                        parts.append(f"…（其余行省略，共 {ws.max_row} 行）")
                        break
                    cells = ["" if v is None else str(v) for v in row]
                    if any(c.strip() for c in cells):
                        parts.append("\t".join(cells))
        finally:
            wb.close()
        return "\n".join(parts)


class PptxExtractor(Extractor):
    name = "pptx"
    exts = (".pptx",)

    def extract(self, path: Path, ctx: ExtractContext) -> str:
        from pptx import Presentation

        try:
            prs = Presentation(str(path))
        except Exception as e:
            raise ExtractError(f"pptx 解析失败: {e}") from e
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"# 第 {i} 页")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        line = "".join(run.text for run in para.runs).strip()
                        if line:
                            parts.append(line)
                if shape.has_table:
                    for row in shape.table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        if any(cells):
                            parts.append("\t".join(cells))
        return "\n".join(parts)
