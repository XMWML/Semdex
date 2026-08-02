from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

from semdex.extractors.office import PptxExtractor


def test_pptx_extractor_includes_table_cells(tmp_path: Path):
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(0.5))
    text_box.text_frame.text = "出差安排"

    table = slide.shapes.add_table(3, 2, Inches(1), Inches(2), Inches(6), Inches(2)).table
    table.cell(0, 0).text = "日期"
    table.cell(0, 1).text = "地点"
    table.cell(1, 0).text = "周一"
    table.cell(1, 1).text = "上海"
    table.cell(2, 0).text = "周二"
    table.cell(2, 1).text = "广州"

    path = tmp_path / "安排.pptx"
    prs.save(path)

    text = PptxExtractor().extract(path, ctx=None)  # type: ignore[arg-type]

    assert text.splitlines() == [
        "# 第 1 页",
        "出差安排",
        "日期\t地点",
        "周一\t上海",
        "周二\t广州",
    ]
